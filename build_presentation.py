"""
presentation.pptx — Social Org Founding Prediction
Executive summary. Business language. Verbal-first structure.
Navy + White theme. One idea per slide. Large text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)
NAVY_MID    = RGBColor(0x1A, 0x35, 0x6E)
BLUE_BRIGHT = RGBColor(0x24, 0x7B, 0xFF)
TEAL        = RGBColor(0x00, 0xC2, 0xB2)
AMBER       = RGBColor(0xFF, 0xB7, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF5, 0xF7, 0xFF)
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xFF)
SLATE       = RGBColor(0x8A, 0x9B, 0xBF)

IMG = "images"
W, H = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=OFF_WHITE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        if line_w:
            shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    kw = dict(width=Inches(w))
    if h:
        kw['height'] = Inches(h)
    slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def dark_slide(slide):
    bg(slide, NAVY)

def light_slide(slide):
    bg(slide, OFF_WHITE)

def top_rule(slide, color=BLUE_BRIGHT, t=1.05):
    box(slide, 0, t, W, 0.045, fill=color)

def slide_title(slide, text, dark_bg=False, size=38, t=0.22, l=0.5, w=12.3):
    color = WHITE if dark_bg else NAVY
    txt(slide, text, l, t, w, 0.85, size=size, bold=True, color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_slide(s)
box(s, 0,    0, W, 0.35, fill=BLUE_BRIGHT)
box(s, 0, 7.15, W, 0.35, fill=TEAL)

txt(s, "Predicting the Founding of\nSocial Organizations",
    0.6, 0.75, 11.8, 2.8,
    size=52, bold=True, color=WHITE)
txt(s, "Among Tecnológico de Monterrey Alumni: A Machine Learning Approach",
    0.6, 3.55, 11.8, 0.65,
    size=22, color=LIGHT_BLUE)

box(s, 0.6, 4.6, 6.5, 0.04, fill=BLUE_BRIGHT)

txt(s, "Gerardo Leiva Díaz  ·  Isaac Sánchez Veloquio\n"
       "David Lozano Arreola  ·  Jorge Betanzo Carriles",
    0.6, 4.8, 10.0, 0.85, size=15, color=LIGHT_BLUE)
txt(s, "Tecnológico de Monterrey  ·  June 2026",
    0.6, 5.75, 10.0, 0.4, size=13, color=SLATE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION TO THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_slide(s)
box(s, 0, 0, W, 0.35, fill=BLUE_BRIGHT)

slide_title(s, "Introduction to the Problem", dark_bg=True, t=0.5, size=38)

txt(s, "Mexico has over 40,000 registered non-profit organizations.",
    0.6, 1.45, 11.8, 0.65, size=22, bold=True, color=WHITE)

txt(s, "Yet we don't know who founds them — or whether a university education\n"
       "in social responsibility actually produces social entrepreneurs.",
    0.6, 2.25, 11.8, 1.0, size=19, color=LIGHT_BLUE)

box(s, 0.6, 3.5, 11.8, 0.05, fill=TEAL)

txt(s, "The question",
    0.6, 3.7, 11.8, 0.45, size=17, bold=True, color=TEAL)
txt(s, "Can we identify, from the characteristics of a Tec alumnus,\n"
       "whether they are likely to have founded a social organization?",
    0.6, 4.25, 11.8, 1.1, size=21, color=WHITE)

txt(s, "QS 80th-anniversary survey  ·  25,356 alumni  ·  Graduating classes 1952–2022",
    0.6, 5.65, 11.8, 0.5, size=15, color=SLATE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR HYPOTHESIS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
light_slide(s)
slide_title(s, "Our Hypothesis")
top_rule(s)

hx = [0.4, 4.55, 8.7]
htitles = ["H1", "H2", "H3"]
hbodies = [
    "Civic behavior\n— volunteering and donations —\nis the strongest predictor\nof founding a social organization.",
    "Accumulated experience\n— age, leadership, prior\nentrepreneurship —\npositively predicts founding.",
    "Tec competency scores\nassessed at graduation\nleave a measurable imprint\nyears later.",
]
for lx, ht, hb in zip(hx, htitles, hbodies):
    box(s, lx, 1.25, 3.9, 4.2, fill=NAVY)
    box(s, lx, 1.25, 3.9, 0.5, fill=BLUE_BRIGHT)
    txt(s, ht, lx + 0.15, 1.29, 3.6, 0.42,
        size=20, bold=True, color=WHITE)
    txt(s, hb, lx + 0.2, 1.88, 3.5, 3.2,
        size=15, color=LIGHT_BLUE)

txt(s, "Three independent hypotheses — all testable with the available data.",
    0.4, 5.7, 12.5, 0.5, size=14, italic=True, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT WE'RE GOING TO DO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
light_slide(s)
slide_title(s, "What We Set Out to Do")
top_rule(s)

steps = [
    ("Predict",
     "Train AI models to identify\nwhich alumni are most likely\nto have founded a social\norganization."),
    ("Explain",
     "Understand which alumni\ncharacteristics drive that\nlikelihood — and in which\ndirection."),
    ("Profile",
     "Describe what the high-\nprobability founder looks like,\nso Tec can identify them\nearly on."),
    ("Recommend",
     "Translate findings into\nconcrete actions for\ninstitutional social impact\nprograms."),
]
for i, (title, body) in enumerate(steps):
    lx = 0.4 + i * 3.2
    box(s, lx, 1.25, 3.0, 5.2, fill=OFF_WHITE, line_color=BLUE_BRIGHT, line_w=1.5)
    box(s, lx, 1.25, 3.0, 0.55, fill=BLUE_BRIGHT)
    txt(s, title, lx + 0.12, 1.28, 2.8, 0.48, size=18, bold=True, color=WHITE)
    txt(s, body,  lx + 0.15, 1.95, 2.75, 4.2, size=14, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — METHODS: AI & PREDICTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
light_slide(s)
slide_title(s, "Method: Prediction with Artificial Intelligence")
top_rule(s)

# Left: what the model does
box(s, 0.4, 1.25, 6.0, 5.5, fill=NAVY)
txt(s, "How does it work?", 0.55, 1.3, 5.7, 0.48,
    size=17, bold=True, color=TEAL)
explanation = [
    ("The model learns patterns",
     "from 25,356 alumni for whom\nwe already know who founded and who didn't."),
    ("It analyzes 64 characteristics",
     "per alumnus: age, career,\nleadership, volunteering, wellbeing…"),
    ("It produces a score",
     "indicating how likely it is\nthat a given alumnus founded\na social organization."),
    ("It is validated on new data",
     "to confirm the learned pattern\nis not just memorization."),
]
y = 2.0
for title, body in explanation:
    txt(s, title, 0.6, y,        5.5, 0.38, size=14, bold=True, color=WHITE)
    txt(s, body,  0.6, y + 0.38, 5.5, 0.65, size=13, color=LIGHT_BLUE, italic=True)
    y += 1.1

# Right: the data snapshot
box(s, 6.7, 1.25, 6.2, 5.5, fill=OFF_WHITE, line_color=LIGHT_BLUE, line_w=1)
txt(s, "The data", 6.85, 1.3, 5.9, 0.48,
    size=17, bold=True, color=NAVY)

img(s, f"{IMG}/target_distribution.png", 6.85, 1.85, 5.8)

txt(s, "Only 7 in every 100 alumni\nfounded a social organization.",
    6.85, 4.65, 5.8, 0.8, size=17, bold=True, color=NAVY)
txt(s, "The model is designed to surface\nthat minority without overlooking them.",
    6.85, 5.55, 5.8, 0.7, size=14, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODELS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
light_slide(s)
slide_title(s, "Models: Six Approaches, One Winner")
top_rule(s)

models = [
    ("Logistic Regression",  "Linear baseline — the simplest approach"),
    ("Decision Tree",        "Interpretable non-linear baseline"),
    ("Random Forest",        "Ensemble of trees — robust to noise"),
    ("Gradient Boosting",    "Sequential learning — best for tabular data"),
    ("Naïve Bayes",          "Density-estimation baseline"),
    ("QDA",                  "Class-specific covariance baseline"),
]
y = 1.3
for name, rationale in models:
    is_winner = name == "Gradient Boosting"
    fill  = BLUE_BRIGHT if is_winner else OFF_WHITE
    lcol  = BLUE_BRIGHT if is_winner else LIGHT_BLUE
    ncol  = WHITE       if is_winner else NAVY
    rcol  = WHITE       if is_winner else SLATE
    box(s, 0.4, y, 12.5, 0.72,
        fill=fill, line_color=lcol, line_w=1.5 if is_winner else 0.8)
    txt(s, name,      0.6,  y + 0.14, 5.5, 0.44, size=16, bold=is_winner, color=ncol)
    txt(s, rationale, 6.3,  y + 0.14, 6.5, 0.44, size=14, color=rcol, italic=True)
    if is_winner:
        txt(s, "WINNER", 11.5, y + 0.14, 1.2, 0.44,
            size=13, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
    y += 0.82

txt(s, "All models were compared on the same data. Gradient Boosting best separated founders from "
       "non-founders — and held up equally well on data it had never seen.",
    0.4, 6.35, 12.5, 0.85, size=14, italic=True, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTS: MODEL PERFORMANCE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
light_slide(s)
slide_title(s, "Results: How Well Does It Work?")
top_rule(s)

img(s, f"{IMG}/only_roc_auc.png", 0.4, 1.2, 7.5)

txt(s, "Gradient Boosting\nis clearly\nahead.",
    8.2, 1.4, 4.9, 1.6, size=32, bold=True, color=NAVY)

box(s, 8.2, 3.15, 4.9, 0.05, fill=BLUE_BRIGHT)

txt(s, "Each curve is a model.\nHigher and further left = better.",
    8.2, 3.35, 4.9, 0.85, size=15, color=SLATE)

box(s, 8.2, 4.35, 4.9, 1.95, fill=NAVY)
txt(s, "By adjusting the model's sensitivity,\nwe go from identifying\n54 → 550 founders\nper 1,000 alumni screened.",
    8.35, 4.45, 4.6, 1.75, size=15, bold=False, color=WHITE)

txt(s, "10× more impact  ·  same model  ·  same data",
    8.2, 6.45, 4.9, 0.5, size=13, color=TEAL, bold=True, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: THE FOUNDER PROFILE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
light_slide(s)
slide_title(s, "Results: Who Is the Founder?")
top_rule(s)

img(s, f"{IMG}/shap_global_importance.png", 0.4, 1.2, 6.8)

# Right: plain-language profile
box(s, 7.4, 1.25, 5.5, 5.6, fill=NAVY)
txt(s, "The founder's profile", 7.55, 1.3, 5.2, 0.48,
    size=16, bold=True, color=TEAL)

traits = [
    ("Active volunteer",               "The strongest signal — by far"),
    ("Has founded another company",    "Prior entrepreneurship as a foundation"),
    ("Older and more experienced",     "Maturity matters more than field"),
    ("Donates to social causes",       "Civic intent expressed through action"),
    ("Has held CEO or board roles",    "Formal leadership as a precursor"),
    ("High entrepreneurship score\nsince university",
                                       "Tec leaves a measurable imprint"),
]
y = 1.95
for trait, note in traits:
    txt(s, trait, 7.6,  y,        5.1, 0.38, size=13, bold=True, color=WHITE)
    txt(s, note,  7.6,  y + 0.36, 5.1, 0.32, size=12, color=LIGHT_BLUE, italic=True)
    box(s, 7.6, y + 0.7, 5.0, 0.02, fill=NAVY_MID)
    y += 0.82

txt(s, "Civic behavior dominates. Demographics alone don't predict — what someone does matters more than who they are.",
    0.4, 6.5, 12.5, 0.75, size=13, italic=True, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSIONS & IMPLICATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_slide(s)
box(s, 0, 0, W, 0.35, fill=TEAL)

slide_title(s, "Conclusions & Implications", dark_bg=True, t=0.5, size=36)

# Left: what we found
box(s, 0.4, 1.4, 5.9, 4.7, fill=NAVY_MID)
txt(s, "What we found", 0.55, 1.45, 5.6, 0.45, size=15, bold=True, color=TEAL)
conclusions = [
    "Volunteering is the most powerful predictor",
    "Experience and leadership matter — not just age",
    "Tec shapes measurable social entrepreneurs decades later",
    "Adjusted sensitivity: 10× more founders identified",
    "The founder profile is clear and actionable",
]
y = 2.05
for c in conclusions:
    txt(s, "▸  " + c, 0.6, y, 5.5, 0.52, size=13, color=LIGHT_BLUE)
    y += 0.56

# Right: what Tec can do
box(s, 6.55, 1.4, 6.35, 4.7, fill=NAVY_MID)
txt(s, "What Tec can do", 6.7, 1.45, 6.1, 0.45, size=15, bold=True, color=AMBER)
actions = [
    ("Identify early",
     "Volunteering · Entrepreneurship competency · Sense of purpose\nAlready measured. No new data needed."),
    ("Prioritize resources",
     "Direct fellowship and incubation programs toward\nprofiles with the highest social impact potential."),
    ("Iterate with data",
     "Re-train the model with new graduating cohorts.\nPilot on one campus before scaling."),
    ("Audit for fairness",
     "The model shows a mild socioeconomic skew —\nreview before any institutional deployment."),
]
y = 2.05
for title, body in actions:
    txt(s, title, 6.75, y,        6.0, 0.38, size=13, bold=True, color=WHITE)
    txt(s, body,  6.75, y + 0.38, 6.0, 0.55, size=12, color=LIGHT_BLUE, italic=True)
    y += 1.05

box(s, 0.4, 6.3, W - 0.8, 0.05, fill=BLUE_BRIGHT)
txt(s, "Three signals already available:  ① Volunteering  ·  ② Entrepreneurship competency  ·  ③ Sense of purpose",
    0.4, 6.45, 12.5, 0.55, size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_slide(s)
box(s, 0, 0, W, 0.35, fill=TEAL)
box(s, 0, 7.15, W, 0.35, fill=BLUE_BRIGHT)

txt(s, "Thank you.",
    0.7, 1.5, 11.8, 2.0, size=72, bold=True, color=WHITE)
txt(s, "Questions & Discussion",
    0.7, 3.5, 11.8, 0.8, size=28, color=LIGHT_BLUE)

box(s, 0.7, 4.55, 8.0, 0.05, fill=BLUE_BRIGHT)

txt(s, "Gerardo Leiva Díaz  ·  Isaac Sánchez Veloquio\n"
       "David Lozano Arreola  ·  Jorge Betanzo Carriles",
    0.7, 4.75, 11.0, 0.85, size=15, color=LIGHT_BLUE)
txt(s, "Tecnológico de Monterrey  ·  June 2026",
    0.7, 5.7, 11.0, 0.4, size=13, color=SLATE, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("presentation.pptx")
size = os.path.getsize("presentation.pptx")
print(f"Saved presentation.pptx  ({size:,} bytes)  —  {len(prs.slides)} slides")
